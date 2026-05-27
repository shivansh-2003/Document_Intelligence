from pptx import Presentation


def extract_full_text(file_path: str) -> str:
    """Extract all text from a PowerPoint as a single string."""
    prs = Presentation(file_path)
    all_text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())
    
    return "\n".join(all_text)


def extract_text_slides(file_path: str, print_slides: bool = False) -> list[dict]:
    """
    Extract text from a PowerPoint slide by slide.
    
    Returns a list of dicts: [{"slide": 1, "text": "..."}, ...]
    Optionally prints each slide with clear formatting.
    """
    prs = Presentation(file_path)
    slides = []
    
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        
        text = "\n".join(slide_texts)
        slide_data = {"slide": i, "text": text}
        slides.append(slide_data)
        
        if print_slides:
            print(f"\n{'='*50}")
            print(f"  SLIDE {i}")
            print(f"{'='*50}")
            print(text if text else "[No text on this slide]")
            print(f"{'─'*50}")
    
    return slides

