"""
table_parser.py
Multi-source table extraction and normalization for RAG pipelines.
Supports: PDF, PPTX, DOCX, Web URLs

Normalizes tabular content into clean markdown with resolved merged cells,
detected headers, extracted footnotes, and quality scoring.
"""

import re
import warnings
from typing import List, Dict, Any, Optional, Union
from dataclasses import dataclass, field
from pathlib import Path

import pandas as pd

# --- Unstructured (primary extractor) ---
from unstructured.partition.auto import partition
from unstructured.partition.pdf import partition_pdf

# --- Web extraction ---
import requests
from bs4 import BeautifulSoup

# --- Optional fallbacks ---
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False


@dataclass
class NormalizedTable:
    """Standardized table object ready for chunking / embedding."""
    source: str
    source_type: str                       # 'pdf', 'pptx', 'docx', 'web'
    page_or_slide: Optional[int]
    table_index: int
    title_or_caption: Optional[str]
    markdown: str
    dataframe: pd.DataFrame
    headers: List[str]
    footnotes: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    quality_score: float = 1.0             # 0.0 - 1.0


class TableParser:
    """
    Extracts and normalizes tables from PDFs, PowerPoint, Word, and web pages.

    Usage:
        parser = TableParser()
        tables = parser.extract("./report.pdf")
        for t in tables:
            print(t.markdown)
    """

    def __init__(
        self,
        use_unstructured: bool = True,
        use_pdfplumber_fallback: bool = True,
        handle_merged_cells: bool = True,
        extract_footnotes: bool = True,
        min_quality: float = 0.3,
    ):
        self.use_unstructured = use_unstructured
        self.use_pdfplumber_fallback = use_pdfplumber_fallback
        self.handle_merged_cells = handle_merged_cells
        self.extract_footnotes = extract_footnotes
        self.min_quality = min_quality

    # =====================================================================
    #  EXTRACTION
    # =====================================================================

    def extract(self, source: Union[str, Path]) -> List[NormalizedTable]:
        """Unified extraction from file path or URL."""
        source_str = str(source)

        if source_str.startswith(("http://", "https://")):
            return self.extract_from_url(source_str)

        path = Path(source)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        suffix = path.suffix.lower()
        if suffix == ".pdf":
            return self.extract_from_pdf(path)
        elif suffix in (".pptx", ".ppt"):
            return self.extract_from_pptx(path)
        elif suffix in (".docx", ".doc"):
            return self.extract_from_docx(path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")

    def extract_from_pdf(self, path: Union[str, Path]) -> List[NormalizedTable]:
        """Extract tables from PDF using unstructured + optional pdfplumber fallback."""
        path = Path(path)
        tables: List[NormalizedTable] = []

        # --- Primary: unstructured ---
        if self.use_unstructured:
            try:
                elements = partition_pdf(
                    str(path),
                    strategy="hi_res",
                    infer_table_structure=True,
                )
                raw_tables = [
                    e for e in elements if e.to_dict().get("type") == "Table"
                ]
                for idx, tbl in enumerate(raw_tables):
                    nt = self._normalize_unstructured_table(
                        tbl,
                        source=str(path),
                        source_type="pdf",
                        table_index=idx,
                        page=None,
                    )
                    if nt.quality_score >= self.min_quality:
                        tables.append(nt)
            except Exception as e:
                warnings.warn(f"Unstructured PDF extraction failed: {e}")

        # --- Fallback: pdfplumber ---
        if not tables and self.use_pdfplumber_fallback and HAS_PDFPLUMBER:
            try:
                with pdfplumber.open(path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        page_tables = page.extract_tables()
                        for tidx, tbl in enumerate(page_tables):
                            nt = self._normalize_raw_grid(
                                tbl,
                                source=str(path),
                                source_type="pdf",
                                table_index=tidx,
                                page=page_num,
                            )
                            if nt.quality_score >= self.min_quality:
                                tables.append(nt)
            except Exception as e:
                warnings.warn(f"pdfplumber fallback failed: {e}")

        return tables

    def extract_from_pptx(self, path: Union[str, Path]) -> List[NormalizedTable]:
        """Extract tables from PowerPoint."""
        path = Path(path)
        tables: List[NormalizedTable] = []

        if self.use_unstructured:
            try:
                elements = partition(filename=str(path))
                raw_tables = [
                    e for e in elements if e.to_dict().get("type") == "Table"
                ]
                for idx, tbl in enumerate(raw_tables):
                    nt = self._normalize_unstructured_table(
                        tbl,
                        source=str(path),
                        source_type="pptx",
                        table_index=idx,
                        page=None,
                    )
                    if nt.quality_score >= self.min_quality:
                        tables.append(nt)
                if tables:
                    return tables
            except Exception as e:
                warnings.warn(f"Unstructured PPTX extraction failed: {e}")

        if HAS_PPTX:
            try:
                prs = Presentation(path)
                for slide_num, slide in enumerate(prs.slides, 1):
                    for shape in slide.shapes:
                        if shape.has_table:
                            grid = []
                            for row in shape.table.rows:
                                grid.append([cell.text for cell in row.cells])
                            nt = self._normalize_raw_grid(
                                grid,
                                source=str(path),
                                source_type="pptx",
                                table_index=len(tables),
                                page=slide_num,
                            )
                            if nt.quality_score >= self.min_quality:
                                tables.append(nt)
            except Exception as e:
                warnings.warn(f"python-pptx fallback failed: {e}")

        return tables

    def extract_from_docx(self, path: Union[str, Path]) -> List[NormalizedTable]:
        """Extract tables from Word documents."""
        path = Path(path)
        tables: List[NormalizedTable] = []

        if self.use_unstructured:
            try:
                elements = partition(filename=str(path))
                raw_tables = [
                    e for e in elements if e.to_dict().get("type") == "Table"
                ]
                for idx, tbl in enumerate(raw_tables):
                    nt = self._normalize_unstructured_table(
                        tbl,
                        source=str(path),
                        source_type="docx",
                        table_index=idx,
                        page=None,
                    )
                    if nt.quality_score >= self.min_quality:
                        tables.append(nt)
                if tables:
                    return tables
            except Exception as e:
                warnings.warn(f"Unstructured DOCX extraction failed: {e}")

        if HAS_DOCX:
            try:
                doc = Document(path)
                for idx, table in enumerate(doc.tables):
                    grid = []
                    for row in table.rows:
                        grid.append([cell.text for cell in row.cells])
                    nt = self._normalize_raw_grid(
                        grid,
                        source=str(path),
                        source_type="docx",
                        table_index=idx,
                        page=None,
                    )
                    if nt.quality_score >= self.min_quality:
                        tables.append(nt)
            except Exception as e:
                warnings.warn(f"python-docx fallback failed: {e}")

        return tables

    def extract_from_url(self, url: str) -> List[NormalizedTable]:
        """Extract tables from web pages."""
        tables: List[NormalizedTable] = []
        try:
            resp = requests.get(
                url, timeout=30, headers={"User-Agent": "Mozilla/5.0"}
            )
            resp.raise_for_status()
            html = resp.text
        except Exception as e:
            warnings.warn(f"Failed to fetch URL {url}: {e}")
            return tables

        # --- Method 1: pandas.read_html ---
        try:
            dfs = pd.read_html(html)
            for idx, df in enumerate(dfs):
                nt = self._normalize_dataframe(
                    df,
                    source=url,
                    source_type="web",
                    table_index=idx,
                    page=None,
                )
                if nt.quality_score >= self.min_quality:
                    tables.append(nt)
        except Exception as e:
            warnings.warn(f"pandas.read_html failed: {e}")

        # --- Method 2: BeautifulSoup for captions & complex tables ---
        if not tables:
            try:
                soup = BeautifulSoup(html, "html.parser")
                for idx, tbl in enumerate(soup.find_all("table")):
                    caption = tbl.find("caption")
                    caption_text = caption.get_text(strip=True) if caption else None
                    rows = []
                    for tr in tbl.find_all("tr"):
                        row = [
                            td.get_text(strip=True)
                            for td in tr.find_all(["td", "th"])
                        ]
                        if row:
                            rows.append(row)
                    if rows:
                        nt = self._normalize_raw_grid(
                            rows,
                            source=url,
                            source_type="web",
                            table_index=idx,
                            page=None,
                        )
                        if caption_text:
                            nt.title_or_caption = caption_text
                        if nt.quality_score >= self.min_quality:
                            tables.append(nt)
            except Exception as e:
                warnings.warn(f"BeautifulSoup fallback failed: {e}")

        return tables

    # =====================================================================
    #  NORMALIZATION
    # =====================================================================

    def _normalize_unstructured_table(
        self,
        element,
        source: str,
        source_type: str,
        table_index: int,
        page: Optional[int],
    ) -> NormalizedTable:
        """Convert an unstructured Table element to NormalizedTable."""
        d = element.to_dict()
        text = d.get("text", "")
        metadata = d.get("metadata", {})

        # unstructured sometimes returns raw text, sometimes HTML-like
        grid = self._text_to_grid(text)
        nt = self._normalize_raw_grid(
            grid,
            source=source,
            source_type=source_type,
            table_index=table_index,
            page=page,
        )
        # Try to recover caption from metadata
        html_meta = metadata.get("text_as_html", "")
        if html_meta and not nt.title_or_caption:
            # Very rough caption extraction
            m = re.search(r"<caption[^>]*>(.*?)</caption>", html_meta, re.I)
            if m:
                nt.title_or_caption = BeautifulSoup(m.group(1), "html.parser").get_text(strip=True)
        return nt

    def _text_to_grid(self, text: str) -> List[List[str]]:
        """Best-effort conversion of unstructured table text to 2D grid."""
        lines = [line.strip() for line in text.split("\n") if line.strip()]
        grid = []
        for line in lines:
            if "|" in line:
                cells = [c.strip() for c in line.split("|") if c.strip()]
            else:
                # Split on 2+ spaces or tabs
                cells = re.split(r"\s{2,}|\t", line)
                cells = [c.strip() for c in cells if c.strip()]
            if cells:
                grid.append(cells)
        return grid

    def _normalize_raw_grid(
        self,
        grid: List[List[str]],
        source: str,
        source_type: str,
        table_index: int,
        page: Optional[int],
        title: Optional[str] = None,
    ) -> NormalizedTable:
        """Normalize a raw 2D list of strings into a clean table."""
        if not grid:
            return NormalizedTable(
                source=source,
                source_type=source_type,
                page_or_slide=page,
                table_index=table_index,
                title_or_caption=title,
                markdown="",
                dataframe=pd.DataFrame(),
                headers=[],
                quality_score=0.0,
            )

        df = pd.DataFrame(grid)

        if self.handle_merged_cells:
            df = self._resolve_merged_cells(df)

        headers, df = self._detect_headers(df)

        footnotes = []
        if self.extract_footnotes:
            df, footnotes = self._extract_footnotes(df)

        df = self._clean_cells(df)
        markdown = self._dataframe_to_markdown(df, headers)
        quality = self._compute_quality(df)

        return NormalizedTable(
            source=source,
            source_type=source_type,
            page_or_slide=page,
            table_index=table_index,
            title_or_caption=title,
            markdown=markdown,
            dataframe=df,
            headers=headers,
            footnotes=footnotes,
            quality_score=quality,
        )

    def _normalize_dataframe(
        self,
        df: pd.DataFrame,
        source: str,
        source_type: str,
        table_index: int,
        page: Optional[int],
    ) -> NormalizedTable:
        """Normalize a pandas DataFrame (e.g. from pd.read_html)."""
        df = df.reset_index(drop=True)

        if self.handle_merged_cells:
            df = self._resolve_merged_cells(df)

        headers, df = self._detect_headers(df)
        df, footnotes = self._extract_footnotes(df)
        df = self._clean_cells(df)
        markdown = self._dataframe_to_markdown(df, headers)
        quality = self._compute_quality(df)

        return NormalizedTable(
            source=source,
            source_type=source_type,
            page_or_slide=page,
            table_index=table_index,
            title_or_caption=None,
            markdown=markdown,
            dataframe=df,
            headers=headers,
            footnotes=footnotes,
            quality_score=quality,
        )

    # -----------------------------------------------------------------
    #  Cell-level normalizers
    # -----------------------------------------------------------------

    def _resolve_merged_cells(self, df: pd.DataFrame) -> pd.DataFrame:
        """Forward-fill empty / NA cells to resolve rowspans & colspans."""
        df = df.copy()
        # Replace empty / whitespace-only with NA, then forward-fill down each column
        df = df.replace(r"^\s*$", pd.NA, regex=True)
        for col in df.columns:
            df[col] = df[col].ffill()
        # For header row, also forward-fill horizontally (colspan simulation)
        if len(df) > 0:
            first_row = df.iloc[0]
            if first_row.isna().any():
                df.iloc[0] = first_row.ffill()
        return df.fillna("")

    def _detect_headers(self, df: pd.DataFrame) -> (List[str], pd.DataFrame):
        """Detect if first row is a header. Returns (headers, df_without_header)."""
        if df.empty:
            return [], df

        first_row = df.iloc[0].astype(str).tolist()

        # Heuristic: contains letters and is reasonably short
        has_letters = any(re.search(r"[a-zA-Z]{2,}", cell) for cell in first_row)
        is_short = all(len(cell) < 100 for cell in first_row)
        looks_like_header = has_letters and is_short

        # Additional heuristic: if row contains typical header words
        header_keywords = [
            "total", "average", "year", "month", "date", "name", "description",
            "amount", "cost", "revenue", "price", "quantity", "id", "no", "number",
            "function", "expenditure", "income", "assets", "liabilities",
        ]
        keyword_score = sum(
            1 for cell in first_row
            if any(kw in cell.lower() for kw in header_keywords)
        )
        if keyword_score >= 2:
            looks_like_header = True

        if looks_like_header:
            headers = [str(h) for h in first_row]
            df = df.iloc[1:].reset_index(drop=True)
        else:
            headers = [f"Col_{i}" for i in range(len(df.columns))]

        # Pad / trim to match column count
        while len(headers) < len(df.columns):
            headers.append(f"Col_{len(headers)}")
        headers = headers[: len(df.columns)]
        df.columns = headers
        return headers, df

    def _extract_footnotes(self, df: pd.DataFrame) -> (pd.DataFrame, List[str]):
        """Extract footnote markers like (1), [1], *, †, ‡ from cells."""
        footnote_pattern = re.compile(r"(\(\d+\)|\[\d+\]|\*|†|‡|§|¶)")
        footnotes = []

        def clean_cell(val):
            if not isinstance(val, str):
                return val
            matches = footnote_pattern.findall(val)
            for m in matches:
                if m not in footnotes:
                    footnotes.append(m)
            return footnote_pattern.sub("", val).strip()

        # pandas >= 2.1 safe apply
        df = df.map(clean_cell)
        return df, footnotes

    def _clean_cells(self, df: pd.DataFrame) -> pd.DataFrame:
        """Strip whitespace, normalize currency, handle accounting negatives (123) → -123."""
        def clean(val):
            if pd.isna(val):
                return ""
            val = str(val).strip()
            val = re.sub(r"\s+", " ", val)          # collapse whitespace
            # Accounting notation: (123) → -123
            if re.match(r"^\(\d[\d,\.]*\)$", val):
                val = "-" + val[1:-1]
            return val

        return df.map(clean)

    def _dataframe_to_markdown(self, df: pd.DataFrame, headers: List[str]) -> str:
        """Convert DataFrame to clean GitHub-flavored markdown table."""
        if df.empty:
            return ""

        lines = []
        lines.append("| " + " | ".join(str(h) for h in headers) + " |")
        lines.append("|" + "|".join(["---"] * len(headers)) + "|")
        for _, row in df.iterrows():
            lines.append("| " + " | ".join(str(v) for v in row.values) + " |")
        return "\n".join(lines)

    def _compute_quality(self, df: pd.DataFrame) -> float:
        """Score 0-1 based on fill-rate and column consistency."""
        if df.empty:
            return 0.0
        total = df.size
        empty = (df == "").sum().sum()
        fill_rate = 1.0 - (empty / total)
        # Penalize wildly inconsistent row lengths (already normalized, so minor)
        row_lengths = df.applymap(len).sum(axis=1)
        consistency = 1.0 if row_lengths.std() < row_lengths.mean() * 0.5 else 0.85
        return round(min(fill_rate * consistency, 1.0), 2)


# =====================================================================
#  HIGH-LEVEL INTEGRATION HELPER
# =====================================================================

def normalize_tables_for_chunking(
    source: Union[str, Path],
    parser: Optional[TableParser] = None,
    enrich_with_llm: bool = False,
    llm_client=None,          # e.g. ChatOllama or OpenAI client
    document_context: str = "",
) -> List[Dict[str, Any]]:
    """
    One-shot function: extract → normalize → optional LLM enrichment → chunk-ready dicts.

    Returns list of dicts with keys:
        - text          : full text for embedding (markdown + caption + footnotes)
        - markdown      : clean markdown only
        - metadata      : source, page, headers, footnotes, quality_score, type
    """
    if parser is None:
        parser = TableParser()

    tables = parser.extract(source)
    chunks: List[Dict[str, Any]] = []

    for tbl in tables:
        parts = []
        if tbl.title_or_caption:
            parts.append(f"Table: {tbl.title_or_caption}")
        parts.append(tbl.markdown)
        if tbl.footnotes:
            parts.append(f"Footnotes: {', '.join(tbl.footnotes)}")

        full_text = "\n\n".join(parts)

        # Optional: LLM contextual enrichment (like the Medium article)
        if enrich_with_llm and llm_client is not None:
            prompt = f"""Given this table and document context, write a 2-sentence 
description of what this table contains. Then keep the markdown table.

Document context:
{document_context[:4000]}

Table:
{tbl.markdown}

Description:"""
            try:
                # Generic invoke pattern
                if hasattr(llm_client, "invoke"):
                    resp = llm_client.invoke(prompt)
                    desc = resp.content if hasattr(resp, "content") else str(resp)
                else:
                    resp = llm_client.chat.completions.create(
                        model="gpt-4o", messages=[{"role": "user", "content": prompt}]
                    )
                    desc = resp.choices[0].message.content
                full_text = f"{desc}\n\n{full_text}"
            except Exception as e:
                warnings.warn(f"LLM enrichment failed for table {tbl.table_index}: {e}")

        chunks.append({
            "text": full_text,
            "markdown": tbl.markdown,
            "source": tbl.source,
            "source_type": tbl.source_type,
            "page": tbl.page_or_slide,
            "table_index": tbl.table_index,
            "headers": tbl.headers,
            "footnotes": tbl.footnotes,
            "quality_score": tbl.quality_score,
            "type": "table",
        })

    return chunks


# =====================================================================
#  EXAMPLE USAGE (uncomment to run)
# =====================================================================
if __name__ == "__main__":
    # --- Single file ---
    # chunks = normalize_tables_for_chunking("./report.pdf")
    # for c in chunks:
    #     print(c["text"][:300], "\n---")

    # --- With LLM enrichment (Ollama example) ---
    # from langchain_ollama import ChatOllama
    # llm = ChatOllama(model="gpt-oss:20b")
    # chunks = normalize_tables_for_chunking(
    #     "./report.pdf",
    #     enrich_with_llm=True,
    #     llm_client=llm,
    #     document_context="Meta Platforms Q2 2024 earnings report",
    # )
    pass